"""Fill Attachment C (Proposal Summary Slide) from official template."""
import shutil
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

C_SRC = (
    "docs/proposals/DARPA_MATHBAC/official_templates/"
    "Attachment_C_Proposal_Summary_Slide_Template.pptx"
)
C_DEST = (
    "docs/proposals/DARPA_MATHBAC/package_work/03_attachment_c_summary_slide/"
    "Attachment_C_SCBE_FILLED.pptx"
)
shutil.copy2(C_SRC, C_DEST)

prs = Presentation(C_DEST)
slide = prs.slides[0]

# ──────────────────────────────────────────────────────────────────────────────
# Content blocks
# ──────────────────────────────────────────────────────────────────────────────

TITLE_TEXT = (
    "SCBE Mathematical Framework for Agentic Communication Protocols\n"
    "ISSAC D DAVIS / SCBE-AETHERMOORE  |  Issac Davis, PI  |  (360) 808-0876  |  "
    "issdandavis7795@gmail.com  |  16 months  |  $839,000  |  UEI: J4NXHM6N5F59  |  "
    "Unclassified / Fundamental Research"
)

TOP_LEFT = (
    "Composed operator T = L₁₄ ∘ ⋯ ∘ L₁ acting on:\n"
    "ψ(t) = (ψ₁(t),…,ψ_N(t), G_t, pd(t)) ∈ Bⁿˣᴺ × G_N × ℝ₊\n\n"
    "Agent state ψ_i(t) embedded in Poincaré ball Bⁿ via L4 (exponential map).\n"
    "Trust-weighted directed Laplacian: L_H(G_t), edge w_ij = H(d_H(ψ_i,ψ_j), pd_ij)\n"
    "Harmonic wall (L12): H(d,pd) = 1/(1 + φ·d_H + 2·pd) ∈ (0,1]\n"
    "CDPTI = Re(λ₂(L_H(G_t))) / Re(λ₂(L_H(G_t0))) — signed coherence drift\n\n"
    "Safety is a geodesic property of T, not a post-hoc filter.\n"
    "BAA alignment: TA1, MC-1 through MC-5"
)

TOP_RIGHT = (
    "PI: Issac Davis — SCBE-AETHERMOORE (sole prop, SAM-active; UEI J4NXHM6N5F59)\n"
    "Sub: Hoags Inc. (bounded DAVA corroborating support, if finalized)\n\n"
    "Benchmark evidence:\n"
    "• Petri adversarial: 172/173 seeds (99.42%); 0.58% false-allow after regex pre-filter\n"
    "• Semantic sphere: 26/26 surfaces, 505 iterations (2026-05-30)\n"
    "• Terminal-bench-core-0.1.1: 13/13 (zero governance overhead)\n"
    "• Real-patch task harness: 5/5 SCBE vs. 0/5 baseline\n\n"
    "Independent derivation: Lyapunov stability, CBF, port-Hamiltonian,\n"
    "persistent excitation, gyroscopic precession (Khalil; Marsden/Ratiu; Cannon)\n\n"
    "Prior government contracts: None  |  SAM: Active"
)

BOTTOM_LEFT = (
    "Phase I Objectives:\n"
    "1. Operator formalization — Lyapunov constants η, b; NMR task families (M3)\n"
    "2. CDPTI ≥ 95% recall, ≤ 5% FPR vs. Mixtral-8x7B baseline (M13)\n"
    "3. Harmonic wall cost curve — φ-weighting ablation; adversarial overhead ratio (M9)\n"
    "4. NMR validation — Task A (1H NMR→structure) + Task B (2D COSY); "
    "ChemBERTa-77M + Qwen2.5-Coder-0.5B; Hammett lane (M9)\n"
    "5. Four PA metrics — MEE, ACV, CDPTI, PIS; IV&V-computable from JSONL bundle\n\n"
    "Why composed operator:\n"
    "Spera (2026, arXiv:2603.15973): 42.6% of 900 real trajectories have conjunctive "
    "multi-agent dependencies — T addresses directly via geometric accumulation.\n"
    "Flat scores, behavioral contracts, post-hoc red-teaming all fail at protocol-graph level."
)

BOTTOM_RIGHT = (
    "Phase I Schedule (16 months, OT §4021):\n"
    "M1 (mo. 1) — Baselines; SSM calibration; IV&V bundle format\n"
    "M3 (mo. 3) — Lyapunov constants; L_H(G_t) instrumented; CDPTI live\n"
    "M6 (mo. 6) — CDPTI demo; ≥2 NMR variants; PI meeting\n"
    "M9 (mo. 9) — Software suite; Hammett generalization; ROMs\n"
    "M13 (mo.13) — Baseline vs. Mixtral; PIS catalog v1; PI meeting\n"
    "M14 (mo.14) — Computational design tool (CDT)\n"
    "M16 (mo.16) — Final report + Phase II plan\n\n"
    "Key metrics:\n"
    "MEE: emission rate (calibrated M3) | ACV: mean ≥0.90 at M9\n"
    "CDPTI: ≥95% recall | PIS: cluster separation at M13\n\n"
    "Cost: $839,000 (42% of $2M cap)  |  Phase II ROM: ~$382,000\n"
    "CAGE: 1EXD5  |  Award: OT for Research (10 U.S.C. §4021)"
)

# ──────────────────────────────────────────────────────────────────────────────
# Helper: replace all text in a text frame, preserving first run formatting
# ──────────────────────────────────────────────────────────────────────────────

def set_tf_text(tf, text, font_size_pt=9):
    """Clear text frame and set content, preserving basic formatting."""
    # Clear existing paragraphs (keep first)
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)

    first_para = tf.paragraphs[0]
    # Clear runs from first paragraph
    for r in list(first_para.runs):
        r._r.getparent().remove(r._r)

    # Split into lines
    lines = text.split("\n")
    first_line = True
    for line in lines:
        if first_line:
            para = first_para
            first_line = False
        else:
            # Add new paragraph
            from pptx.oxml.ns import qn
            from lxml import etree
            new_p = copy.deepcopy(first_para._p)
            # Clear runs in the copy
            for r in new_p.findall(qn('a:r')):
                new_p.remove(r)
            tf._txBody.append(new_p)
            para = tf.paragraphs[-1]

        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_size_pt)
        run.font.name = "Tahoma"
        run.font.color.rgb = RGBColor(0, 0, 0)

    tf.word_wrap = True


# ──────────────────────────────────────────────────────────────────────────────
# Fill shapes by their template text (identity)
# ──────────────────────────────────────────────────────────────────────────────

SHAPE_MAP = {
    "Proposal Title": TITLE_TEXT,       # title placeholder — first line only
    "Provide graphic.": TOP_LEFT,
    "Describe new ideas.": TOP_RIGHT,
    "Describe the risks of the effort.": BOTTOM_LEFT,
    "Describe the team and responsibilities.": BOTTOM_RIGHT,
}

filled = set()
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    current = shape.text_frame.text.strip()
    # Match by checking if any key appears in current text
    for key, replacement in SHAPE_MAP.items():
        if key in current:
            if key == "Proposal Title":
                # Title placeholder: set via runs to preserve formatting
                set_tf_text(shape.text_frame, replacement, font_size_pt=7)
            else:
                set_tf_text(shape.text_frame, replacement, font_size_pt=9)
            print(f"  FILLED [{key[:40]!r}]")
            filled.add(key)
            break

print(f"\nFilled {len(filled)}/{len(SHAPE_MAP)} zones.")
if len(filled) < len(SHAPE_MAP):
    missing = set(SHAPE_MAP) - filled
    print(f"  WARNING: unfilled zones: {missing}")

prs.save(C_DEST)
print(f"Saved: {C_DEST}")
