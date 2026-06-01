"""Fill Attachment X (Proposal Overview and Proposed Metrics) from official template."""
import shutil
import copy
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

X_SRC = (
    "docs/proposals/DARPA_MATHBAC/official_templates/"
    "Attachment_X_Proposal_Overview_and_Proposed_Metrics.pptx"
)
X_DEST = (
    "docs/proposals/DARPA_MATHBAC/package_work/10_attachment_x_metrics/"
    "Attachment_X_SCBE_FILLED.pptx"
)
shutil.copy2(X_SRC, X_DEST)
prs = Presentation(X_DEST)


# ──────────────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────────────

def delete_slide(prs, idx):
    """Remove slide at index idx from presentation."""
    sldIdLst = prs.slides._sldIdLst
    rId = sldIdLst[idx].get(qn("r:id"))
    prs.part.drop_rel(rId)
    del sldIdLst[idx]


def set_tf(tf, text, pt=9, bold_first_line=False):
    """Replace all text in a text frame."""
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)
    first_para = tf.paragraphs[0]
    for r in list(first_para.runs):
        r._r.getparent().remove(r._r)

    lines = text.split("\n")
    first_line = True
    for line in lines:
        if first_line:
            para = first_para
            first_line = False
        else:
            new_p = copy.deepcopy(first_para._p)
            for r in new_p.findall(qn("a:r")):
                new_p.remove(r)
            tf._txBody.append(new_p)
            para = tf.paragraphs[-1]
        run = para.add_run()
        run.text = line
        run.font.size = Pt(pt)
        run.font.name = "Tahoma"
        run.font.color.rgb = RGBColor(0, 0, 0)
    tf.word_wrap = True


def set_cell(table, row, col, text, pt=9, bold=False):
    """Set a table cell's text."""
    cell = table.cell(row, col)
    tf = cell.text_frame
    if tf.paragraphs:
        for r in list(tf.paragraphs[0].runs):
            r._r.getparent().remove(r._r)
        while len(tf.paragraphs) > 1:
            tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(pt)
    run.font.name = "Tahoma"
    run.font.color.rgb = RGBColor(0, 0, 0)
    run.font.bold = bold
    tf.word_wrap = True


def clear_shape(shape):
    """Clear all text from a shape's text frame."""
    if shape.has_text_frame:
        set_tf(shape.text_frame, "", pt=9)


# ──────────────────────────────────────────────────────────────────────────────
# Step 1: Delete guidance slide (index 0) and TA2 slide (index 2 → becomes 1 after deletion)
# ──────────────────────────────────────────────────────────────────────────────

print(f"Before deletion: {len(prs.slides)} slides")
delete_slide(prs, 0)   # guidance slide
delete_slide(prs, 1)   # TA2 slide (was index 2, now index 1)
print(f"After deletion: {len(prs.slides)} slides")
# Now: slide 0 = TA1 overview, slide 1 = metric template


# ──────────────────────────────────────────────────────────────────────────────
# Step 2: Fill TA1 overview slide (index 0)
# ──────────────────────────────────────────────────────────────────────────────

ta1 = prs.slides[0]

for shape in ta1.shapes:
    txt = shape.text_frame.text.strip() if shape.has_text_frame else ""
    name = shape.name

    # Platform description text box
    if "# of agents:" in txt:
        set_tf(shape.text_frame, (
            "# of agents: 2 (configurable per campaign)\n"
            "types of agents: Math-reasoning orchestrator + SSM science validator\n"
            "models: ChemBERTa-77M + Qwen2.5-Coder-0.5B-Instruct\n"
            "# tasks / domains: 2 NMR task families + 1 generalization lane (Hammett)\n"
            "months to build (current) platform: 18"
        ), pt=9)
        print(f"  FILLED: platform description")

    # TA1 Rubric table (4 rows x 3 cols: header + 3 data rows)
    if shape.has_table and shape.table.rows[0].cells[0].text.strip() == "TA1 Rubric":
        t = shape.table
        # Headers already: TA1 Rubric | Current results | Month 4 Target
        set_cell(t, 1, 0, "% success rate")
        set_cell(t, 1, 1, "Terminal-bench: 13/13 (100%)\nReal-patch: 5/5 (100%)\nPetri: 172/173 (99.42%)")
        set_cell(t, 1, 2, "CDPTI ≥95% recall vs. Mixtral; ACV_norm ≥0.90")
        set_cell(t, 2, 0, "Governance overhead")
        set_cell(t, 2, 1, "Zero overhead on neutral tasks (terminal-bench-core-0.1.1)")
        set_cell(t, 2, 2, "Maintain ≤5% overhead on governed benchmark runs")
        set_cell(t, 3, 0, "Benchmark protocol")
        set_cell(t, 3, 1, "IV&V JSONL bundle format defined; SSM probe-layer TBD")
        set_cell(t, 3, 2, "IV&V bundle locked at M1; φ-ablation complete by M3")
        print(f"  FILLED: TA1 rubric table")

    # TA1 Goals table (8 rows x 3 cols)
    if shape.has_table and shape.table.rows[0].cells[0].text.strip().startswith("TA1 Goal"):
        t = shape.table
        rows = [
            ("(a)", "Poincaré-ball operator construction (L1–L14)",
             "Hyperbolic geometry (Ganea, Gülçehre), Lyapunov stability (Khalil), composed operators (Mac Lane)"),
            ("(b)", "Trust-weighted directed Laplacian L_H(G_t)",
             "Spectral graph theory (Newman), signed graph dynamics (Strogatz), Fiedler connectivity"),
            ("(c)", "CDPTI signed coherence drift index",
             "Spera non-comp. safety (arXiv:2603.15973), persistent homology (Cohen-Steiner)"),
            ("(d)", "MEE artifact extraction & verification",
             "Minerva (arXiv:2206.14858), LeanDojo (NeurIPS 2023), AlphaGeometry (Nature 2024)"),
            ("(e)", "ACV 5-predicate structural compliance vector",
             "Pearl causality, von Neumann unitarity, Mac Lane composition, Goyal RIM"),
            ("(f)", "PIS protocol identity signature",
             "Cover/Thomas info theory, Tishby bottleneck, Indyk/Motwani ANN, persistent homology"),
            ("(g)", "CDT: Lagrangian automation for new subdomains",
             "Phase I deliverable C; NMR → Hammett generalization lane validated by M9"),
        ]
        for i, (goal, method, ref) in enumerate(rows):
            row_idx = i + 1
            if row_idx < len(t.rows):
                set_cell(t, row_idx, 0, goal)
                set_cell(t, row_idx, 1, method)
                set_cell(t, row_idx, 2, ref)
        print(f"  FILLED: TA1 goals table")

    # TA1 Details table (5 rows x 2 cols)
    if shape.has_table and shape.table.rows[0].cells[0].text.strip().startswith("TA1"):
        t = shape.table
        if len(t.rows) >= 4 and len(t.columns) == 2:
            set_cell(t, 0, 0, "TA1 Field", bold=True)
            set_cell(t, 0, 1, "Details", bold=True)
            set_cell(t, 1, 0, "Selected Science Subdomain(s)")
            set_cell(t, 1, 1, "NMR spectroscopy (primary); cheminformatics / Hammett equation (generalization)")
            set_cell(t, 2, 0, "1st family of scientific tasks")
            set_cell(t, 2, 1, "Task A: 1H NMR spectrum → molecular structure (BMRB+PDB)")
            set_cell(t, 3, 0, "2nd family of scientific tasks")
            set_cell(t, 3, 1, "Task B: 2D COSY → vicinal J-coupling connectivity (BMRB)")
            if len(t.rows) > 4:
                set_cell(t, 4, 0, "Small science models (SSMs)")
                set_cell(t, 4, 1, "ChemBERTa-77M (primary); Qwen2.5-Coder-0.5B-Instruct (reasoning)")
            print(f"  FILLED: TA1 details table")

    # Clear instruction text boxes
    if "Instructions (delete before submitting)" in txt:
        clear_shape(shape)
        print(f"  CLEARED: instruction box")
    if "How did you set your targets" in txt:
        clear_shape(shape)
    if "Include image of current agentic platform" in txt:
        clear_shape(shape)


# ──────────────────────────────────────────────────────────────────────────────
# Metric content for 4 slides
# ──────────────────────────────────────────────────────────────────────────────

METRICS = [
    {
        "title": "Proposed Metric #1 — Mathematical Evidence Emission (MEE)",
        "content": (
            "Technical Area: TA1\n\n"
            "Description: MEE measures the rate and verification quality of checkable mathematical "
            "artifacts emitted during inter-agent communication: formulas, derivation steps, citations, "
            "executable code, and numerical assertions with units.\n\n"
            "Method for calculation: For each act, extract evidence artifacts and score each with a "
            "class-appropriate verifier (symbolic simplification, sandbox execution, dimensional checks, "
            "citation resolution). MEE = sum(artifact verification scores) / number of inter-agent "
            "messages. Report MEE_density, MEE_pass_rate, and artifact-class balance. IV&V inputs: "
            "governance_events.jsonl, evidence_artifacts.jsonl, task_manifests.json.\n\n"
            "Progress against program goals: MATHBAC seeks science-discovery protocols that are "
            "understandable, not only successful. MEE distinguishes a correct black-box answer from a "
            "campaign that leaves a mathematical audit trail.\n\n"
            "Supplement to current metrics: Program success, speed, and generalization are end-state "
            "measures. MEE is a trace-level measure of whether the protocol produced falsifiable "
            "scientific evidence during the campaign.\n\n"
            "Adoption argument: MEE gives IV&V a low-cost audit surface: the evidence log is the same "
            "material needed to check whether agent communication contained science-relevant mathematical "
            "content. Vendor-neutral; computable from JSONL without SCBE source.\n\n"
            "Literature: Minerva (arXiv:2206.14858), Polu/Sutskever (arXiv:2009.03393), "
            "AlphaGeometry (Nature 2024), LeanDojo (NeurIPS 2023).\n\n"
            "Phase I targets: M3 — extractor and verifier calibration. M9 — MEE density/pass-rate "
            "on NMR task families. M13 — compare MEE distribution vs. Mixtral-8x7B baseline."
        ),
    },
    {
        "title": "Proposed Metric #2 — Axiom Compliance Vector (ACV)",
        "content": (
            "Technical Area: TA1\n\n"
            "Description: ACV is a five-component vector in [0,1]^5 measuring per-act compliance with "
            "Unitarity, Locality, Causality, Symmetry, and Composition.\n\n"
            "Method for calculation: For each act, evaluate five predicates: information conservation in "
            "handoff (Unitarity), bounded influence radius (Locality), dependency-graph topological order "
            "(Causality), invariance under role permutation tests (Symmetry), pipeline schema continuity "
            "(Composition). Campaign ACV = mean predicate vector. "
            "ACV_norm = (5 - ||1 − ACV||_1) / 5. IV&V inputs: governance_events.jsonl, "
            "latent_states.npz, task_manifests.json.\n\n"
            "Progress against program goals: TA1 asks for stability and convergence analysis. ACV "
            "measures whether the protocol maintains the structural conditions under which the composed "
            "operator is well-posed.\n\n"
            "Supplement to current metrics: A protocol can be fast and correct while violating causality, "
            "symmetry, or composition. ACV catches structurally pathological protocols that end-state "
            "metrics cannot distinguish.\n\n"
            "Adoption argument: ACV is vendor-neutral — any performer can emit the five predicate values "
            "from ordinary communication and latent-state logs. Also provides Phase II teams a "
            "compatibility check before combining protocols.\n\n"
            "Literature: von Neumann operator formalism, Pearl causality, Mac Lane composition, "
            "Goyal recurrent independent mechanisms, symmetry-invariance literature.\n\n"
            "Phase I targets: M3 — ACV field added to event schema. M9 — mean ACV_norm ≥ 0.90 on "
            "calibrated NMR campaigns (threshold subject to IV&V baseline adjustment)."
        ),
    },
    {
        "title": "Proposed Metric #3 — Communication-Dynamics Phase-Transition Index (CDPTI)",
        "content": (
            "Technical Area: TA1\n\n"
            "Description: CDPTI measures signed regime changes in the multi-agent communication graph. "
            "It detects both communication breakdown and coordinated adversarial convergence.\n\n"
            "Method for calculation: Build time-windowed directed graph G_t. Edge (i,j) weighted by "
            "communication energy, signed by harmonic-wall evaluation: "
            "w_ij(t) = sign(H_ij(t) − H_min) × ||log(S_ij)||. Compute trust-weighted directed Laplacian "
            "L_H(G_t) = D_out(W_H) − W_H. Signed readout: "
            "CDPTI(t) = Re(λ₂(L_H(G_t))) / Re(λ₂(L_H(G_t0))). "
            "Negative signed-connectivity growth triggers escalation even when unsigned connectivity rises. "
            "IV&V inputs: governance_events.jsonl, latent_states.npz, protocol_graph_edges.jsonl.\n\n"
            "Progress against program goals: PA line 443 asks performers to characterize and quantify "
            "progress of communication dynamics. CDPTI identifies when the protocol changes regime and "
            "whether that change points toward the safe basin or an unsafe attractor.\n\n"
            "Supplement to current metrics: Outcome metrics do not show whether a campaign progressed "
            "smoothly, stalled, or converged adversarially. CDPTI provides temporal and directional "
            "resolution.\n\n"
            "Adoption argument: CDPTI directly addresses non-compositional safety risks identified by "
            "Spera (2026): conjunctive multi-agent drift appears as signed graph motion.\n\n"
            "Literature: Newman graph spectra, Strogatz nonlinear dynamics, Cohen-Steiner persistence "
            "stability, Spera non-compositional safety (arXiv:2603.15973).\n\n"
            "Phase I targets: M3 — instrument L_H(G_t). M6 — live CDPTI computation. "
            "M13 — injected-drift recall ≥95% and FPR ≤5%."
        ),
    },
    {
        "title": "Proposed Metric #4 — Protocol Information Signature (PIS)",
        "content": (
            "Technical Area: TA1\n\n"
            "Description: PIS is a fixed-dimensional protocol fingerprint summarizing the structural "
            "identity of a multi-agent communication protocol independent of agent names and surface "
            "formatting.\n\n"
            "Method for calculation: Compute mutual-information features, conditional-entropy graph "
            "features, role-transition features, latent-delta covariance summaries, and persistent-homology "
            "features of the protocol graph. Project the resulting feature vector into a fixed-dimensional "
            "signature using the released pis_projection_matrix.json. Compare protocols with cosine "
            "distance and cluster persistence. IV&V inputs: governance_events.jsonl, latent_states.npz, "
            "protocol_graph_edges.jsonl, pis_projection_matrix.json.\n\n"
            "Progress against program goals: TA1 deliverable C requires a catalog of protocol design "
            "principles and domain-specific optimal protocols. PIS is the catalog key: lets IV&V determine "
            "whether two protocols are structurally similar or genuinely distinct.\n\n"
            "Supplement to current metrics: Success and speed do not identify whether a protocol is new, "
            "redundant, or transferable. PIS adds identity and similarity structure to the rubric.\n\n"
            "Adoption argument: PIS helps detect benchmark gaming, supports cross-subdomain transfer "
            "analysis, and allows protocol catalog entries to be queried by structural similarity.\n\n"
            "Literature: Cover/Thomas information theory, Tishby information bottleneck, Indyk/Motwani "
            "ANN, Charikar similarity estimation, persistent-homology stability.\n\n"
            "Phase I targets: M6 — preliminary PIS feature extractor. "
            "M13 — PIS catalog v1 with high-performance reference centroids and cluster-separation report."
        ),
    },
]


# ──────────────────────────────────────────────────────────────────────────────
# Step 3: Fill and duplicate metric slides
# ──────────────────────────────────────────────────────────────────────────────

def fill_metric_slide(slide, metric):
    """Fill a metric slide with content from metric dict."""
    for shape in slide.shapes:
        txt = shape.text_frame.text.strip() if shape.has_text_frame else ""

        # Title
        if shape.has_text_frame and "Proposed Metric" in txt and shape.name.startswith("Title"):
            set_tf(shape.text_frame, metric["title"], pt=16)
            continue

        # Main content box
        if "Technical Area:" in txt or "Description:" in txt:
            set_tf(shape.text_frame, metric["content"], pt=9)
            continue

        # Delete instruction box
        if "Each proposed metric is limited to one slide" in txt:
            clear_shape(shape)
            continue


def duplicate_slide_from_template(prs, template_slide):
    """Duplicate a slide, returning the new slide."""
    layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    # Remove default placeholder shapes added by add_slide
    for shape in list(new_slide.placeholders):
        sp = shape._element
        sp.getparent().remove(sp)
    # Copy all shapes from template
    for shape in template_slide.shapes:
        el = shape._element
        new_slide.shapes._spTree.insert(2, copy.deepcopy(el))
    return new_slide


# Metric template is now slide index 1
metric_template = prs.slides[1]

# Fill metric 1 in the template slide
fill_metric_slide(metric_template, METRICS[0])
print(f"  FILLED: {METRICS[0]['title'][:50]}")

# Add slides 2-4 by duplicating template
for m in METRICS[1:]:
    new_slide = duplicate_slide_from_template(prs, metric_template)
    fill_metric_slide(new_slide, m)
    print(f"  ADDED+FILLED: {m['title'][:50]}")

prs.save(X_DEST)
print(f"\nSaved: {X_DEST} ({len(prs.slides)} slides total)")
