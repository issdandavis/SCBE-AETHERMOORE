"""Inspect slide/placeholder structure of Attachment C and X pptx templates."""
from pptx import Presentation
from pptx.util import Pt
import sys

def inspect(path, label):
    print(f"\n{'='*60}")
    print(f"  {label}")
    print(f"  {path}")
    print('='*60)
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"  ERROR loading: {e}")
        return

    print(f"  Slides: {len(prs.slides)}")
    print(f"  Slide size: {prs.slide_width.pt:.0f}pt x {prs.slide_height.pt:.0f}pt")

    for si, slide in enumerate(prs.slides):
        print(f"\n  --- Slide {si+1} ---")
        for i, shape in enumerate(slide.shapes):
            shtype = str(shape.shape_type)
            name = shape.name
            has_tf = shape.has_text_frame
            left = shape.left.pt if shape.left else 0
            top = shape.top.pt if shape.top else 0
            w = shape.width.pt if shape.width else 0
            h = shape.height.pt if shape.height else 0

            print(f"  [{i}] '{name}' type={shtype} pos=({left:.0f},{top:.0f}) size={w:.0f}x{h:.0f} has_tf={has_tf}")

            if has_tf:
                tf = shape.text_frame
                full_text = tf.text.strip()
                print(f"       TEXT: {full_text[:200]!r}")
                for pi, para in enumerate(tf.paragraphs[:5]):
                    for ri, run in enumerate(para.runs[:3]):
                        rt = run.text or ""
                        if rt.strip():
                            print(f"         para[{pi}] run[{ri}]: {rt[:100]!r}")

            # Check if it's a placeholder
            try:
                if shape.is_placeholder:
                    ph = shape.placeholder_format
                    print(f"       PLACEHOLDER idx={ph.idx} type={ph.type}")
            except Exception:
                pass

            # Tables
            if shape.has_table:
                t = shape.table
                print(f"       TABLE {t.rows.__len__()}x{t.columns.__len__()}")
                for r in range(min(3, len(t.rows))):
                    row_text = [t.cell(r, c).text[:30] for c in range(len(t.columns))]
                    print(f"         row[{r}]: {row_text}")

inspect(
    "docs/proposals/DARPA_MATHBAC/official_templates/Attachment_C_Proposal_Summary_Slide_Template.pptx",
    "ATTACHMENT C — Proposal Summary Slide"
)

inspect(
    "docs/proposals/DARPA_MATHBAC/official_templates/Attachment_X_Proposal_Overview_and_Proposed_Metrics.pptx",
    "ATTACHMENT X — Proposal Overview and Proposed Metrics"
)
